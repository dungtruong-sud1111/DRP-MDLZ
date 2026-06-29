import sys
import os
from pptx import Presentation

# Configure stdout to UTF-8
sys.stdout.reconfigure(encoding='utf-8')

def audit_presentation(filepath="deck.pptx"):
    if not os.path.exists(filepath):
        print(f"ERROR: File {filepath} not found.")
        sys.exit(1)
        
    print(f"\n=======================================================")
    print(f"RUNNING SELF-HEALING AUDIT ON: {filepath}")
    print(f"=======================================================")
    
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"CRITICAL ERROR: Failed to load presentation. {e}")
        sys.exit(1)
        
    violations_count = 0
    
    # 1. Check Slide Count & Mandatory Agenda Slide (Rule 4)
    if len(prs.slides) < 2:
        print("WARNING: Slide deck has less than 2 slides. Agenda slide cannot be verified.")
    else:
        # Check if Slide 2 is Agenda
        agenda_slide = prs.slides[1]
        has_agenda_title = False
        for shape in agenda_slide.shapes:
            if shape.has_text_frame and "AGENDA" in shape.text_frame.text.upper():
                has_agenda_title = True
                break
        if not has_agenda_title:
            print("❌ Rule 4 VIOLATION (Mandatory Agenda): Slide 2 is not an Agenda slide or lacks 'AGENDA' text.")
            violations_count += 1
            
    # Trace through slides
    for idx, slide in enumerate(prs.slides, 1):
        slide_title = "Unknown Slide"
        for s in slide.shapes:
            if s.has_text_frame and s.text_frame.text:
                first_line = s.text_frame.text.split("\n")[0]
                if len(first_line) > 3 and "Page" not in first_line:
                    slide_title = first_line[:40]
                    break
                    
        print(f"\nAnalyzing Slide {idx}: '{slide_title}'")
        
        # Collect shape text boxes and bounding boxes for overlap audit
        text_shapes = []
        
        for sh_idx, shape in enumerate(slide.shapes):
            # A. Safe Vertical Boundary (Rule 1)
            btm = (shape.top + shape.height) / 914400 # in inches
            if btm > 7.15:
                # Exclude footer branding and page number
                t_content = ""
                if shape.has_text_frame:
                    t_content = shape.text_frame.text.strip().replace("\n", " ")
                
                if "Page" not in t_content and len(t_content) > 0:
                    print(f"  ❌ Rule 1 VIOLATION (Overflow): Shape boundary at {btm:.2f} in > 7.15 in. Content: '{t_content[:30]}...'")
                    violations_count += 1
            
            # Save standard textbox coordinates for overlap checks
            if shape.has_text_frame and shape.text_frame.text.strip():
                t = shape.text_frame.text.strip()
                if "Page" not in t: # Ignore footers
                    # Left, Top, Right, Bottom in EMU
                    bbox = (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)
                    text_shapes.append((shape, bbox, t))
                    
            # B. Font Uniformity (Rule 6)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name and run.font.name != "Be Vietnam Pro":
                            print(f"  ❌ Rule 6 VIOLATION (Font): Found non-standard font '{run.font.name}' in paragraph: '{paragraph.text[:20]}'")
                            violations_count += 1
                            
            # C. Color Contrast (Rule 5)
            if shape.has_text_frame and shape.fill and shape.fill.type == 1: # Solid fill
                bg_color = shape.fill.fore_color.rgb if shape.fill.fore_color.type == 1 else None
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.color and run.font.color.type == 1:
                            fg_color = run.font.color.rgb
                            # Check if both are same color
                            if bg_color and fg_color and bg_color == fg_color:
                                print(f"  ❌ Rule 5 VIOLATION (Contrast): Invisible text (foreground matches background color) in: '{run.text[:20]}'")
                                violations_count += 1
                                
        # D. Overlaps & Collisions (Rule 2)
        for i in range(len(text_shapes)):
            for j in range(i+1, len(text_shapes)):
                sh_a, bbox_a, text_a = text_shapes[i]
                sh_b, bbox_b, text_b = text_shapes[j]
                
                # Check intersection bounding box
                # bbox format: (left, top, right, bottom)
                overlap_x = min(bbox_a[2], bbox_b[2]) - max(bbox_a[0], bbox_b[0])
                overlap_y = min(bbox_a[3], bbox_b[3]) - max(bbox_a[1], bbox_b[1])
                
                if overlap_x > 180000 and overlap_y > 180000: # Overlap more than 0.2 inch
                    # Check nesting: if shape B is completely inside shape A (e.g. text in card)
                    nested_a = (bbox_a[0] <= bbox_b[0] + 18000 and bbox_a[1] <= bbox_b[1] + 18000 and 
                                bbox_a[2] + 18000 >= bbox_b[2] and bbox_a[3] + 18000 >= bbox_b[3])
                    nested_b = (bbox_b[0] <= bbox_a[0] + 18000 and bbox_b[1] <= bbox_a[1] + 18000 and 
                                bbox_b[2] + 18000 >= bbox_a[2] and bbox_b[3] + 18000 >= bbox_a[3])
                    
                    if not (nested_a or nested_b):
                        print(f"  ❌ Rule 2 VIOLATION (Overlap): Collision detected!")
                        print(f"     Shape A: '{text_a[:30].replace(chr(10), ' ')}'")
                        print(f"     Shape B: '{text_b[:30].replace(chr(10), ' ')}'")
                        violations_count += 1
                        
    print(f"\n=======================================================")
    if violations_count == 0:
        print("🎉 AUDIT PASSED! Slide deck is 100% clean and professional.")
        print("=======================================================")
        return True
    else:
        print(f"❌ AUDIT FAILED! Found {violations_count} styling violations.")
        print("   Please check the logs above and correct build_deck.py coordinates.")
        print("=======================================================")
        return False

if __name__ == "__main__":
    audit_presentation()
